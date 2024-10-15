# -*- coding: utf-8 -*-
# Copyright (c) Huawei Technologies Co., Ltd. 2024. All rights reserved.
import itertools

from loguru import logger
from paddleocr import PaddleOCR
from pptx import Presentation
from langchain_core.documents import Document
from langchain_community.document_loaders.base import BaseLoader

from mx_rag.document.loader.base_loader import BaseLoader as mxBaseLoader
from mx_rag.utils.file_check import SecFileCheck, FileCheckError, PathNotFileException


class PowerPointLoader(BaseLoader, mxBaseLoader):
    EXTENSION = (".pptx",)
    MAX_SIZE = 100 * 1024 * 1024
    MAX_TABLE_ROW = 100
    MAX_TABLE_COL = 50

    def __init__(self, file_path, lang="ch"):
        super().__init__(file_path)
        try:
            self.ocr = PaddleOCR(use_angle_cls=True, lang=lang)
        except Exception as err:
            raise ValueError(f"init ocr failed, {err}") from err

    def lazy_load(self):
        self._check_file_valid()
        try:
            return self._load_ppt()
        except ValueError as ve:
            logger.error(f"Value error: {ve}")
            return iter([])
        except Exception as err:
            logger.error(f"load '{self.file_path}' failed, {err}")
            return iter([])

    def _check_file_valid(self):
        SecFileCheck(self.file_path, self.MAX_SIZE).check()
        if not self.file_path.endswith(self.EXTENSION):
            raise TypeError("file type not correct")
        if self._is_zip_bomb():
            raise ValueError(f"'{self.file_path}' is a risk of zip bombs")

    def _load_merged_cell(self, data, cell, row, col):
        span_height = cell.span_height
        span_width = cell.span_width
        for span_row in range(row, row + span_height):
            for span_col in range(col, col + span_width):
                data[span_row][span_col] = cell.text

    def _load_table(self, table):
        rows = min(len(table.rows), self.MAX_TABLE_ROW)
        cols = min(len(table.columns), self.MAX_TABLE_COL)
        # 初始化一个二维列表来存储表格数据
        data = [["" for _ in range(cols)] for _ in range(rows)]
        for row in range(rows):
            for col in range(cols):
                cell = table.cell(row, col)
                if not cell.text:
                    continue
                if not cell.is_merge_origin:
                    data[row][col] = cell.text
                    continue
                self._load_merged_cell(data, cell, row, col)

        return itertools.chain.from_iterable(data)

    def _load_image_text(self, image_bytes):
        result = self.ocr.ocr(image_bytes, cls=True)
        try:
            res = [line[1][0] for line in result[0]]
            return res
        except TypeError as err:
            logger.info(f"can not load text from image, {err}")
            return None

    def _load_slide(self, slide):
        slide_text = []
        for shape in slide.shapes:
            # 识别图片中的文字
            if hasattr(shape, "image"):
                image_data = shape.image.blob
                img_text = self._load_image_text(image_data)
                if img_text is not None:
                    slide_text.extend(img_text)

            # 检查形状是否为表格
            if shape.has_table:
                table = shape.table
                table_text = self._load_table(table)
                slide_text.extend(table_text)

            # 获取AutoShape中的纯文本
            if shape.has_text_frame:
                slide_text.append(shape.text_frame.text.replace("\n", " "))
        return " ".join(slide_text)

    def _load_ppt(self):
        prs = Presentation(self.file_path)
        for slide in prs.slides:
            slide_text = self._load_slide(slide)
            yield Document(page_content=slide_text, metadata={"source": self.file_path})
